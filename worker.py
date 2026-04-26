# =============================================================================
# SECTION 1: IMPORTS & INITIAL SETUP
# =============================================================================
# All standard library, third-party, and project-specific imports are declared
# here at the top so every dependency is visible in one place.

import os
import asyncio
import hashlib
import json
import re
import time
from collections import defaultdict
from typing import List, Tuple, Dict, Any
from urllib.parse import urlparse, parse_qsl, urlencode
from pptx import Presentation
import base64
import pandas as pd

# Third-party libraries
import aiohttp
import fitz  # PyMuPDF
import nltk
import openai
import redis
from bs4 import BeautifulSoup
from docx import Document
from dotenv import load_dotenv
from langchain_text_splitters import RecursiveCharacterTextSplitter
from sentence_transformers import CrossEncoder
from astrapy import DataAPIClient
from langdetect import detect
from deep_translator import GoogleTranslator

# Perform initial setup tasks
print(" Worker starting up...")
load_dotenv()
nltk.download('punkt', quiet=True)
nltk.download('stopwords', quiet=True)
os.environ["TOKENIZERS_PARALLELISM"] = "false"


# =============================================================================
# SECTION 2: CONFIGURATION & INITIALIZATION
# =============================================================================
# All tunable constants, environment-variable loaders, and service clients live
# here. Keeping them together means changing a model name or weight only ever
# requires touching this one section.

# --- Constants ---
EMBEDDING_MODEL_API = "text-embedding-3-small"
EMBEDDING_DIMENSION = 1536
VECTOR_COLLECTION_NAME = "ragdb_semantic"
KEYWORD_INDEX_COLLECTION_NAME = "ragdb_keyword"
PROCESSED_DOCS_KEY = "v2" # Redis key for tracking processed docs

def load_environment_config() -> Dict[str, str]:
    """Loads required environment variables and returns them as a dictionary."""
    config = {
        "redis_url": os.getenv("REDIS_URL"),
        "openai_api_key": os.getenv("OPENAI_API_KEY"),
        "astra_db_api_endpoint": os.getenv("ASTRA_DB_API_ENDPOINT"),
        "astra_db_application_token": os.getenv("ASTRA_DB_APPLICATION_TOKEN"),
    }
    if not all(config.values()):
        raise ValueError("One or more essential environment variables are missing.")
    return config

def initialize_services(config: Dict[str, str]) -> Dict[str, Any]:
    """Initializes and returns clients for external services and ML models."""
    print("️ Initializing external services and models...")

    # Initialize clients
    redis_conn = redis.from_url(config["redis_url"])
    sync_openai_client = openai.OpenAI(api_key=config["openai_api_key"])
    async_openai_client = openai.AsyncOpenAI(api_key=config["openai_api_key"])
    
    print(" Loading CrossEncoder model for reranking...")
    reranker_model = CrossEncoder('cross-encoder/ms-marco-MiniLM-L-6-v2')
    print(" CrossEncoder model loaded.")
    
    print("️ Connecting to Astra DB...")
    astra_client = DataAPIClient(config["astra_db_application_token"])
    db = astra_client.get_database(config["astra_db_api_endpoint"])
    print(f" Connected to Astra DB. Namespace: {db.namespace}")

    # Setup Astra DB collections
    collection_names = db.list_collection_names()
    
    # Vector collection
    if VECTOR_COLLECTION_NAME in collection_names:
        print(f"Collection '{VECTOR_COLLECTION_NAME}' already exists. Loading it.")
        vector_collection = db.get_collection(VECTOR_COLLECTION_NAME)
    else:
        print(f"Collection '{VECTOR_COLLECTION_NAME}' not found. Creating it...")
        vector_collection = db.create_collection(
            VECTOR_COLLECTION_NAME,
            dimension=EMBEDDING_DIMENSION
        )
        print(" Vector collection created.")

    # Keyword index collection
    if KEYWORD_INDEX_COLLECTION_NAME in collection_names:
        print(f"Collection '{KEYWORD_INDEX_COLLECTION_NAME}' already exists. Loading it.")
        keyword_index_collection = db.get_collection(KEYWORD_INDEX_COLLECTION_NAME)
    else:
        print(f"Collection '{KEYWORD_INDEX_COLLECTION_NAME}' not found. Creating it...")
        keyword_index_collection = db.create_collection(KEYWORD_INDEX_COLLECTION_NAME)
        print(" Keyword index collection created.")

    print(" Astra DB Collections ready.")
    print("-" * 50)
    
    return {
        "redis_conn": redis_conn,
        "sync_openai_client": sync_openai_client,
        "async_openai_client": async_openai_client,
        "reranker_model": reranker_model,
        "db": db,
        "vector_collection": vector_collection,
        "keyword_index_collection": keyword_index_collection
    }

# --- RAG Pipeline Configuration ---
# Central place for all tunable parameters of the RAG pipeline.
RAG_CONFIG = {
    "semantic_weight": 0.7,
    "keyword_weight": 0.3,
    "k_semantic": 30,
    "k_keyword": 30,
    "k_rerank": 15
}
print(f" RAG pipeline configured with default tuning parameters: {RAG_CONFIG}")


# =============================================================================
# SECTION 3: UTILITY FUNCTIONS
# =============================================================================
# Lightweight, stateless helper functions used across the entire pipeline.
# These have no external dependencies and are the easiest functions to test
# in isolation.

def logtime(phase: str, start: float, last: float = None) -> float:
    """Logs the time elapsed since the start and the last checkpoint."""
    now = time.perf_counter()
    since_start = now - start
    log_message = f"️ [{phase}] {since_start:.2f}s"
    if last:
        since_last = now - last
        log_message += f" (+{since_last:.2f}s since last)"
    print(log_message)
    return now

def generate_document_id(source: str) -> str:
    """
    Generates a stable, unique document ID from a URL or file path.
    Hashes the final ID to prevent excessive length.

    For URLs, it normalises the path and query string before hashing so that
    two requests for the same page always produce the same ID.
    For file paths, the basename is used as the base of the ID.
    """
    if urlparse(source).scheme in ['http', 'https']:
        parsed = urlparse(source)
        base_id = parsed.path.strip("/").replace("/", "-") or "root"
        if parsed.query:
            query_string = urlencode(sorted(parse_qsl(parsed.query)))
            full_id = f"{base_id}?{query_string}"
        else:
            full_id = base_id
    else: # It's a file path
        full_id = os.path.basename(source)
    
    hashed = hashlib.md5(full_id.encode()).hexdigest()[:8]
    sanitized_base = re.sub(r'[^a-zA-Z0-9_-]', '', os.path.splitext(full_id)[0])
    return f"{sanitized_base}-{hashed}"


# =============================================================================
# SECTION 4: TEXT EXTRACTION — MULTI-FORMAT DOCUMENT PARSERS
# =============================================================================
# Each function in this group is responsible for extracting raw text from one
# specific file format (PDF, PPTX, DOCX, HTML, Excel, Image).
# They all follow the same contract: accept a file path (or content), return
# a (text, page_count) tuple.  This uniformity makes it easy to add new
# formats later without changing downstream code.

def extract_text_from_pdf(file_path: str) -> Tuple[str, int]:
    """Extracts all text and page count from a PDF file."""
    doc = fitz.open(file_path)
    full_text = "".join(page.get_text() for page in doc)
    num_pages = doc.page_count
    doc.close()
    return full_text, num_pages

def extract_text_from_pptx(file_path: str) -> Tuple[str, int]:
    """Extracts all text and slide count from a PPTX file."""
    prs = Presentation(file_path)
    full_text = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                full_text.append(shape.text)
    return "\n".join(full_text), len(prs.slides)

def extract_text_from_docx(file_path: str) -> Tuple[str, int]:
    """Extracts all text and estimates page count from a DOCX file."""
    doc = Document(file_path)
    full_text = "\n".join([para.text for para in doc.paragraphs if para.text.strip()])
    para_count = len(doc.paragraphs)
    num_pages = max(1, para_count // 5)  # Heuristic for page count
    return full_text, num_pages

def extract_text_from_html(html_content: str) -> Tuple[str, int]:
    """Extracts clean text and estimates page count from HTML content."""
    soup = BeautifulSoup(html_content, "html.parser")
    for tag in soup(["script", "style", "noscript", "header", "footer", "nav"]):
        tag.extract()
    text = soup.get_text(separator="\n", strip=True)
    num_pages = max(1, len(text.split()) // 300) # Heuristic for page count
    return text, num_pages

async def scrape_and_clean_url(session: aiohttp.ClientSession, url: str) -> str:
    """Asynchronously scrapes and cleans text content from a given URL."""
    try:
        async with session.get(url, timeout=15) as response:
            response.raise_for_status()
            html = await response.text()
            clean_text, _ = extract_text_from_html(html)
            return clean_text
    except Exception as e:
        print(f"️ Could not scrape URL {url}: {e}")
        return ""
    
def extract_text_from_excel(file_path: str) -> Tuple[str, int]:
    """Extracts text from all sheets of an Excel file."""
    xls = pd.ExcelFile(file_path)
    full_text = []
    for sheet_name in xls.sheet_names:
        df = pd.read_excel(xls, sheet_name=sheet_name, header=None)
        # Add a header for context and convert sheet to string
        sheet_text = f"--- Sheet: {sheet_name} ---\n{df.to_string(index=False)}"
        full_text.append(sheet_text)
    return "\n\n".join(full_text), len(xls.sheet_names)

async def extract_text_from_image_async(file_path: str, services: Dict[str, Any]) -> Tuple[str, int]:
    """Extracts text from an image using OpenAI's vision model."""
    print("️ Extracting text from image using Vision API...")
    with open(file_path, "rb") as image_file:
        base64_image = base64.b64encode(image_file.read()).decode('utf-8')

    response = await services["async_openai_client"].chat.completions.create(
        model="gpt-4o-mini",
        messages=[
            {
                "role": "user",
                "content": [
                    {"type": "text", "text": "Extract all text from this image exactly as it appears. If no text is present, say nothing."},
                    {"type": "image_url", "image_url": {"url": f"data:image/jpeg;base64,{base64_image}"}},
                ],
            }
        ],
        max_tokens=2048,
    )
    extracted_text = response.choices[0].message.content
    return extracted_text, 1 # An image is considered a single page


# =============================================================================
# SECTION 5: DOCUMENT ANALYSIS — CLASSIFICATION & CHUNKING
# =============================================================================
# Before text can be indexed it needs to be split into manageable pieces
# ("chunks"). This section groups the logic that decides *how* to split:
#   1. Classify the document type from its content.
#   2. Pick appropriate chunk size / overlap for that type.
#   3. Actually split the text using LangChain's splitter.
#
# Keeping classification and chunking together makes it easy to tune one
# document type without affecting the others.

def simple_doc_type_classifier(text: str) -> str:
    """
    Classifies document type based on keyword heuristics for chunking strategy.

    Scans the lowercased text for domain-specific trigger words and returns a
    category string ('policy', 'scientific', 'narrative', or 'general').
    The returned category is consumed by get_chunking_params to select the
    best chunk size for that document type.
    """
    lowered = text.lower()
    if any(word in lowered for word in ["policy", "contract", "agreement", "legal"]):
        return "policy"
    elif any(word in lowered for word in ["research", "study", "technical", "scientific"]):
        return "scientific"
    elif any(word in lowered for word in ["story", "novel", "narrative", "literature"]):
        return "narrative"
    return "general"

def get_chunking_params(doc_type: str, num_pages: int) -> Tuple[int, int]:
    """
    Determines chunk size and overlap based on document type and length.

    Returns a (chunk_size, chunk_overlap) tuple that controls how the text
    splitter divides the document.  Both values can be tuned per doc_type
    without touching any other part of the pipeline.
    """
    # This logic is kept as per the original script. It can be expanded.
    chunk_size = 1000
    chunk_overlap = 100
    # print(f"ℹ️ Using chunking params: size={chunk_size}, overlap={chunk_overlap} for type='{doc_type}'")
    return chunk_size, chunk_overlap

def chunk_text(text: str, doc_type: str, num_pages: int) -> List[str]:
    """
    Splits text into chunks using a recursive character splitter.

    Delegates parameter selection to get_chunking_params so that the
    right strategy is automatically applied for each document type.
    """
    chunk_size, chunk_overlap = get_chunking_params(doc_type, num_pages)
    splitter = RecursiveCharacterTextSplitter(
        chunk_size=chunk_size, 
        chunk_overlap=chunk_overlap
    )
    return splitter.split_text(text)


# =============================================================================
# SECTION 6: KEYWORD INDEX — TOKENISATION & INVERTED INDEX CONSTRUCTION
# =============================================================================
# This section implements the *keyword* half of the hybrid retrieval system.
# An inverted index maps every meaningful token in the corpus to the set of
# chunk IDs that contain it — the same idea used by traditional search engines.
#
# Steps:
#   1. NLTK stopwords are loaded once at module level.
#   2. tokenize_for_indexing() cleans and filters a text string.
#   3. build_and_store_inverted_index() builds the mapping and writes it to
#      the Astra DB keyword collection in batches to avoid timeout errors.

import nltk
from nltk.corpus import stopwords

# Create a set of standard English stopwords from the NLTK library.
# The try/except block will automatically download the list if it's not present.
try:
    NLTK_STOPWORDS = set(stopwords.words('english'))
except LookupError:
    print("Downloading NLTK stopwords list...")
    nltk.download('stopwords')
    NLTK_STOPWORDS = set(stopwords.words('english'))

print(f" Standard NLTK stopword list initialized with {len(NLTK_STOPWORDS)} words.")


def tokenize_for_indexing(text: str) -> List[str]:
    """
    Lowercases, tokenizes text, and removes standard NLTK stopwords for
    efficient and relevant keyword indexing.

    Words like 'the', 'is', and 'and' are removed so only semantically
    meaningful tokens are stored in the inverted index.
    """
    tokens = re.findall(r'\w+', text.lower())
    # Filter out any token that is in our NLTK_STOPWORDS set
    return [token for token in tokens if token not in NLTK_STOPWORDS]


def build_and_store_inverted_index(
    document_id: str, chunks: List[str], keyword_collection
):
    """
    Builds an inverted index from document chunks and stores it in Astra DB.

    For every chunk, each meaningful token is mapped to the chunk's ID.
    The resulting index documents are written to the keyword collection
    in batches of 200 to avoid Astra DB's 10-second request timeout.
    """
    inverted_index = defaultdict(set)
    for i, chunk in enumerate(chunks):
        chunk_id = f"{document_id}-chunk-{i}"
        for token in tokenize_for_indexing(chunk):
            inverted_index[token].add(chunk_id)

    index_docs = [
        {
            "_id": f"{document_id}-{token}",
            "document_id": document_id,
            "token": token,
            "chunks": list(chunk_ids)
        }
        for token, chunk_ids in inverted_index.items()
    ]
    
    if index_docs:
        # FIX: Send to Astra DB in batches of 200 to prevent 10-second timeout
        batch_size = 200
        for i in range(0, len(index_docs), batch_size):
            keyword_collection.insert_many(index_docs[i : i + batch_size])
            
    print(f" Inverted index stored for {len(inverted_index)} unique tokens.")


# =============================================================================
# SECTION 7: VECTOR EMBEDDING GENERATION
# =============================================================================
# Dense vector embeddings turn raw text into numerical representations that
# capture semantic meaning.  Similar sentences end up with similar vectors,
# which is what makes semantic (vector) search work.
#
# Two variants are provided:
#   • get_embeddings_sync  — used at query time inside a worker thread.
#   • get_embeddings_batch_async — used at index time to embed many chunks
#     concurrently without blocking the event loop.

def get_embeddings_sync(texts: List[str], client: openai.OpenAI) -> List[List[float]]:
    """
    Synchronously generates embeddings for a list of texts.

    Used in blocking contexts (e.g. inside asyncio.to_thread) where an
    async call would not be safe.  Newlines are stripped before sending
    because they can degrade embedding quality.
    """
    response = client.embeddings.create(
        input=[t.replace('\n', ' ') for t in texts], 
        model=EMBEDDING_MODEL_API
    )
    return [e.embedding for e in response.data]

async def get_embeddings_batch_async(
    texts: List[str], client: openai.AsyncOpenAI, batch_size: int = 500
) -> List[List[float]]:
    """
    Generates embeddings for a list of texts in batches asynchronously.

    Splits the input into batches of `batch_size` and fires all requests
    concurrently using asyncio.gather, then re-assembles the results in
    the original order.  This is significantly faster than sequential calls
    when indexing large documents with hundreds of chunks.
    """
    all_embeddings = []
    tasks = []
    for i in range(0, len(texts), batch_size):
        batch = [t.replace('\n', ' ') for t in texts[i:i + batch_size]]
        tasks.append(client.embeddings.create(input=batch, model=EMBEDDING_MODEL_API))
    
    responses = await asyncio.gather(*tasks)
    for response in responses:
        all_embeddings.extend([e.embedding for e in response.data])
    return all_embeddings


# =============================================================================
# SECTION 8: LANGUAGE DETECTION & TRANSLATION
# =============================================================================
# To support multilingual documents, the pipeline detects the source language
# and, if it differs from English, appends a machine translation so that
# English queries can still match non-English content during retrieval.

async def translate_if_needed_async(text: str, target_lang: str = "en") -> str:
    """
    Detects language and translates it to the target language if different,
    using the deep-translator library.

    Runs both detection and translation in a thread pool (via asyncio.to_thread)
    so the async event loop is never blocked by these CPU/network-bound calls.
    Returns an empty string when no translation is needed or if the process fails.
    """
    try:
        # Step 1: Detect the language of the input text.
        detected_lang = await asyncio.to_thread(detect, text)

        # Step 2: Check if translation is necessary.
        if detected_lang != target_lang:
            print(f" Detected language: {detected_lang}, translating to {target_lang}...")

            # Step 3: Perform the translation using the correct deep-translator syntax.
            translated_text = await asyncio.to_thread(
                GoogleTranslator(source='auto', target=target_lang).translate,
                text
            )

            # Step 4: Return the formatted, translated text.
            if translated_text:
                return f"\n\n--- AUTO TRANSLATION ({detected_lang} → {target_lang}) ---\n{translated_text}\n--- END TRANSLATION ---"
            else:
                return ""  # Return empty string if translation result is empty.
        else:
            # The text is already in the target language.
            return ""

    except Exception as e:
        print(f"️ Language detection/translation failed: {e}")
        return ""


# =============================================================================
# SECTION 9: CORE INDEXING ENGINE
# =============================================================================
# _process_and_index_document is the shared backbone called by every
# file-type pipeline.  It receives extracted text and orchestrates the four
# indexing steps in order:
#   1. Classify & chunk the text.
#   2. Build & persist the keyword inverted index.
#   3. Generate dense vector embeddings asynchronously.
#   4. Insert chunk vectors into Astra DB and mark the document as done in Redis.
#
# Centralising these steps here means a bug fix or optimisation applies to
# every supported file format automatically.

async def _process_and_index_document(
    document_id: str,
    text: str,
    num_pages: int,
    services: Dict[str, Any]
):
    """
    Internal function to handle the common processing steps and mark as
    processed in Redis upon success.

    This function is called by every format-specific pipeline (PDF, DOCX, URL,
    etc.) after text extraction.  It should not be called directly from outside
    this module — use the appropriate process_and_index_* function instead.
    """
    t0 = time.perf_counter()
    
    # 1. Chunk the text
    doc_type = simple_doc_type_classifier(text)
    chunks = chunk_text(text, doc_type, num_pages)
    if not chunks:
        print(f"️ No chunks could be generated for '{document_id}'. Aborting.")
        return
    t1 = logtime(f"Split text into {len(chunks)} chunks", t0)

    # 2. Build and store the inverted index
    build_and_store_inverted_index(document_id, chunks, services["keyword_index_collection"])
    t2 = logtime("Built inverted keyword index", t0, t1)
    
    # 3. Generate dense vector embeddings
    print(f"Embedding {len(chunks)} chunks...")
    dense_embeddings = await get_embeddings_batch_async(chunks, services["async_openai_client"])
    t3 = logtime("Generated dense embeddings", t0, t2)
    
    # 4. Insert documents into vector collection
    documents_to_insert = [
        {
            "_id": f"{document_id}-chunk-{i}", "document_id": document_id,
            "text": chunk, "$vector": dense,
        }
        for i, (chunk, dense) in enumerate(zip(chunks, dense_embeddings))
    ]
    
    # FIX: Send vectors to Astra DB in batches of 50 to prevent timeouts
    batch_size = 50
    for i in range(0, len(documents_to_insert), batch_size):
        services["vector_collection"].insert_many(documents_to_insert[i : i + batch_size])
        
    t4 = logtime("Inserted all chunks into Astra DB", t0, t3)
    
    # 5. Mark as processed in Redis on success
    services["redis_conn"].sadd(PROCESSED_DOCS_KEY, document_id)
    print(f" Marked '{document_id}' as processed in Redis.")
    
    print(f" Finished indexing '{document_id}' successfully.")
    print("-" * 50)


# =============================================================================
# SECTION 10: FILE PROCESSING PIPELINES (PER FORMAT, WITH REDIS DEDUP CHECK)
# =============================================================================
# Each function here is the public entry point for a specific file type.
# They all follow the same three-step pattern:
#   1. Check Redis — if the document was already indexed, skip it immediately.
#   2. Extract raw text using the appropriate parser from Section 4.
#   3. Hand off to _process_and_index_document (Section 9) for the rest.
#
# This deduplication check prevents re-indexing the same file on repeated runs,
# saving both time and vector storage quota.

async def process_and_index_pdf(file_path: str, document_id: str, services: Dict[str, Any]):
    """
    Orchestrates the full processing pipeline for a PDF file, using a
    PRE-GENERATED document_id.

    Extra step vs other formats: any URLs found in the PDF body are also
    scraped and their content appended to the text before indexing.
    """
    if services["redis_conn"].sismember(PROCESSED_DOCS_KEY, document_id):
        print(f"ℹ️ Document ID '{document_id}' (from PDF) already processed. Skipping.")
        return

    print(f"️ Starting pipeline for PDF: {file_path} (ID: {document_id})")
    full_text, num_pages = extract_text_from_pdf(file_path)
    if not full_text.strip():
        print(f"️ No text could be extracted from the PDF '{file_path}'. Aborting.")
        return

    found_urls = sorted(list(set(re.findall(r'https?://[^\s<>"]+|www\.[^\s<>"]+', full_text))))
    combined_text = full_text

    if found_urls:
        print(f" Found {len(found_urls)} unique URLs. Scraping...")
        async with aiohttp.ClientSession() as session:
            scraping_tasks = [scrape_and_clean_url(session, url) for url in found_urls]
            scraped_results = await asyncio.gather(*scraping_tasks)
            scraped_texts = [
                f"--- LINKED CONTENT FROM: {url} ---\n{text}\n--- END LINKED CONTENT ---)"
                for url, text in zip(found_urls, scraped_results) if text
            ]
            if scraped_texts:
                combined_text += "\n\n" + "\n\n".join(scraped_texts)
                print(f" Appended content from {len(scraped_texts)} URLs.")

    #  Automatic language detection & translation
    translation_context = await translate_if_needed_async(combined_text)
    if translation_context:
        combined_text += translation_context

    await _process_and_index_document(document_id, combined_text, num_pages, services)


async def process_and_index_docx(file_path: str, document_id: str, services: Dict[str, Any]):
    """
    Orchestrates the full processing pipeline for a DOCX file, using a
    PRE-GENERATED document_id.
    """
    if services["redis_conn"].sismember(PROCESSED_DOCS_KEY, document_id):
        print(f"ℹ️ Document ID '{document_id}' (from DOCX) already processed. Skipping.")
        return

    print(f"️ Starting pipeline for DOCX: {file_path} (ID: {document_id})")
    full_text, num_pages = extract_text_from_docx(file_path)
    if not full_text.strip():
        print(f"️ No text could be extracted from the DOCX '{file_path}'. Aborting.")
        return
    await _process_and_index_document(document_id, full_text, num_pages, services)


async def process_and_index_url(url: str, document_id: str, services: Dict[str, Any]):
    """
    Orchestrates the full processing pipeline for a web URL, using a
    PRE-GENERATED document_id.
    """
    if services["redis_conn"].sismember(PROCESSED_DOCS_KEY, document_id):
        print(f"ℹ️ Document ID '{document_id}' (from URL) already processed. Skipping.")
        return

    print(f"️ Starting pipeline for URL: {url} (ID: {document_id})")
    async with aiohttp.ClientSession() as session:
        html_content = await scrape_and_clean_url(session, url)
    if not html_content:
        print(f"️ Could not retrieve or parse content from URL '{url}'. Aborting.")
        return
        
    full_text, num_pages = extract_text_from_html(html_content)
    if not full_text.strip():
        print(f"️ No text content found after cleaning HTML from '{url}'. Aborting.")
        return
    await _process_and_index_document(document_id, full_text, num_pages, services)


async def process_and_index_pptx(file_path: str, document_id: str, services: Dict[str, Any]):
    """Orchestrates the full processing pipeline for a PPTX file."""
    if services["redis_conn"].sismember(PROCESSED_DOCS_KEY, document_id):
        print(f"ℹ️ Document ID '{document_id}' (from PPTX) already processed. Skipping.")
        return
    
    print(f"️ Starting pipeline for PPTX: {file_path} (ID: {document_id})")
    full_text, num_pages = extract_text_from_pptx(file_path)
    if not full_text.strip():
        print(f"️ No text could be extracted from the PPTX '{file_path}'. Aborting.")
        return
    await _process_and_index_document(document_id, full_text, num_pages, services)


async def process_and_index_excel(file_path: str, document_id: str, services: Dict[str, Any]):
    """Orchestrates the full processing pipeline for an Excel file."""
    if services["redis_conn"].sismember(PROCESSED_DOCS_KEY, document_id):
        print(f"ℹ️ Document ID '{document_id}' (from Excel) already processed. Skipping.")
        return
        
    print(f"️ Starting pipeline for Excel: {file_path} (ID: {document_id})")
    full_text, num_sheets = extract_text_from_excel(file_path)
    if not full_text.strip():
        print(f"️ No text could be extracted from the Excel file '{file_path}'. Aborting.")
        return
    await _process_and_index_document(document_id, full_text, num_sheets, services)

async def process_and_index_image(file_path: str, document_id: str, services: Dict[str, Any]):
    """Orchestrates the full processing pipeline for an image file."""
    if services["redis_conn"].sismember(PROCESSED_DOCS_KEY, document_id):
        print(f"ℹ️ Document ID '{document_id}' (from Image) already processed. Skipping.")
        return

    print(f"️ Starting pipeline for Image: {file_path} (ID: {document_id})")
    full_text, num_pages = await extract_text_from_image_async(file_path, services)
    if not full_text.strip():
        print(f"️ No text could be extracted from the Image file '{file_path}'. Aborting.")
        return
    await _process_and_index_document(document_id, full_text, num_pages, services)

async def process_and_index_zip(file_path: str, document_id: str, services: Dict[str, Any]):
    """
    Orchestrates the processing of a ZIP archive by processing its contents.

    Extracts the archive to a temporary directory, then iterates over each
    contained file, dispatching to the correct extractor based on file extension.
    All extracted texts are combined and indexed as a single logical document.
    """
    if services["redis_conn"].sismember(PROCESSED_DOCS_KEY, document_id):
        print(f"ℹ️ Document ID '{document_id}' (from ZIP) already processed. Skipping.")
        return

    print(f"️ Starting pipeline for ZIP: {file_path} (ID: {document_id})")
    all_texts = []
    total_pages = 0

    with tempfile.TemporaryDirectory() as temp_dir:
        with zipfile.ZipFile(file_path, 'r') as zip_ref:
            zip_ref.extractall(temp_dir)
        
        print(f"️ Extracted {len(os.listdir(temp_dir))} files from zip archive.")
        for filename in os.listdir(temp_dir):
            extracted_file_path = os.path.join(temp_dir, filename)
            file_ext = os.path.splitext(filename)[-1].lower()
            text, pages = "", 0
            
            try:
                print(f"   -> Processing '{filename}'...")
                if file_ext == '.pdf':
                    text, pages = extract_text_from_pdf(extracted_file_path)
                elif file_ext == '.docx':
                    text, pages = extract_text_from_docx(extracted_file_path)
                elif file_ext == '.pptx':
                    text, pages = extract_text_from_pptx(extracted_file_path)
                elif file_ext in ['.xlsx', '.xls']:
                    text, pages = extract_text_from_excel(extracted_file_path)
                elif file_ext in ['.jpg', '.jpeg', '.png']:
                    # Image extraction is async, so we await it here.
                    text, pages = await extract_text_from_image_async(extracted_file_path, services)
                else:
                    print(f"   -> ️ Unsupported file type '{file_ext}' in ZIP. Skipping.")
                    continue
                
                if text.strip():
                    all_texts.append(f"--- START CONTENT FROM: {filename} ---\n\n{text}\n\n--- END CONTENT FROM: {filename} ---")
                    total_pages += pages
            except Exception as e:
                print(f"   ->  Error processing '{filename}' inside ZIP: {e}")

    if not all_texts:
        print(f"️ No text could be extracted from any files in the ZIP '{file_path}'. Aborting.")
        return

    combined_text = "\n\n".join(all_texts)
    await _process_and_index_document(document_id, combined_text, total_pages, services)


# =============================================================================
# SECTION 11: HYBRID RETRIEVAL — SEMANTIC & KEYWORD SEARCH
# =============================================================================
# This is the core of the *hybrid* retrieval strategy.  Given a user query,
# the system searches the corpus using two complementary methods in parallel:
#
#   • Semantic search  — embeds the query and finds nearest vectors in Astra DB.
#                        Catches paraphrase-style matches even without shared words.
#   • Keyword search   — tokenizes the query, looks up each token in the
#                        inverted index, and ranks chunks by term co-occurrence.
#                        Catches exact-match and rare-term queries that semantic
#                        search sometimes misses.
#
# The two result sets are merged (deduped), then optionally reranked by a
# CrossEncoder model that scores each (query, chunk) pair more precisely.

def _blocking_semantic_search(
    query: str, document_id: str, services: Dict[str, Any], k_semantic: int
) -> List[str]:
    """
    Performs a synchronous semantic vector search against Astra DB.

    Encodes the query into a dense vector and retrieves the top-k most similar
    chunks filtered to the specified document.  Runs synchronously so it can
    be safely offloaded to a thread via asyncio.to_thread.
    """
    start_time = time.perf_counter()
    dense_embedding = get_embeddings_sync([query], services["sync_openai_client"])[0]
    t1 = logtime(f"[Semantic-Thread] Query Encoding", start_time)
    results = services["vector_collection"].find(
        sort={"$vector": dense_embedding}, limit=k_semantic, projection={"text": 1}, filter={"document_id": document_id}
    )
    contexts = [doc['text'] for doc in results if 'text' in doc]
    logtime(f"[Semantic-Thread] Astra Vector Search ({len(contexts)} candidates)", start_time, t1)
    return contexts

def _blocking_keyword_search(
    query: str, document_id: str, services: Dict[str, Any], k_keyword: int
) -> List[str]:
    """
    Performs an intelligent, ranked, and limited keyword search.

    Tokenizes the query, looks up matching entries in the inverted index, then
    ranks candidate chunks by how many query tokens they contain (term
    frequency voting).  Runs synchronously for the same reason as the semantic
    search above.
    """
    start_time = time.perf_counter()
    query_tokens = tokenize_for_indexing(query)
    if not query_tokens:
        logtime("[Keyword-Thread] No keywords left after stopword removal", start_time)
        return []
    
    keyword_ids = [f"{document_id}-{token}" for token in query_tokens]
    keyword_results = services["keyword_index_collection"].find(filter={"_id": {"$in": keyword_ids}})
    t1 = logtime(f"[Keyword-Thread] Inverted Index Lookup for {len(query_tokens)} keywords", start_time)

    chunk_scores = defaultdict(int)
    for doc in keyword_results:
        for chunk_id in doc.get("chunks", []):
            chunk_scores[chunk_id] += 1
            
    if not chunk_scores:
        logtime("[Keyword-Thread] No chunks found for keywords", start_time, t1)
        return []

    sorted_chunks = sorted(chunk_scores.items(), key=lambda item: item[1], reverse=True)
    top_chunk_ids = [chunk_id for chunk_id, score in sorted_chunks[:k_keyword]]
    
    context_docs = services["vector_collection"].find(filter={"_id": {"$in": top_chunk_ids}}, projection={"text": 1})
    
    text_map = {doc['_id']: doc['text'] for doc in context_docs}
    contexts = [text_map[chunk_id] for chunk_id in top_chunk_ids if chunk_id in text_map]
    
    logtime(f"[Keyword-Thread] Ranked & Retrieved top {len(contexts)} chunks", start_time, t1)
    return contexts

async def hybrid_retrieve_async(
    original_query: str, document_id: str, services: Dict[str, Any], k_semantic: int, k_keyword: int
) -> List[Tuple[str, str]]:
    """
    Asynchronously performs hybrid retrieval by running semantic and keyword
    searches in parallel.

    Both blocking search functions are dispatched to the thread pool via
    asyncio.to_thread so they run concurrently.  Results are merged into a
    single deduplicated list with each chunk tagged by its source
    ('semantic', 'keyword', or 'both') for use by the weighted reranker.
    """
    print(f"\n Starting parallel retrieval for query on '{document_id}'...")
    start_time = time.perf_counter()
    semantic_task = asyncio.to_thread(_blocking_semantic_search, original_query, document_id, services, k_semantic)
    keyword_task = asyncio.to_thread(_blocking_keyword_search, original_query, document_id, services, k_keyword)
    semantic_contexts, keyword_contexts = await asyncio.gather(semantic_task, keyword_task)
    t1 = logtime("[Main-Thread] Parallel retrieval tasks finished", start_time)

    context_source_map = {}
    for context in keyword_contexts: context_source_map[context] = 'keyword'
    for context in semantic_contexts:
        context_source_map[context] = 'both' if context in context_source_map else 'semantic'
    combined_results = list(context_source_map.items())
    
    print(f" Semantic: {len(semantic_contexts)}, Keyword: {len(keyword_contexts)}, Combined unique: {len(combined_results)}")
    logtime(f"[Main-Thread] Total retrieval for '{document_id}'", start_time, t1)
    return combined_results


# =============================================================================
# SECTION 12: RERANKING — CROSS-ENCODER WEIGHTED SCORING
# =============================================================================
# Retrieval returns many candidate chunks quickly but imprecisely.  Reranking
# re-scores every (query, chunk) pair using a heavier CrossEncoder model that
# reads *both* texts together, producing a much more accurate relevance score.
#
# The system also applies weights based on how a chunk was retrieved:
#   • 'semantic' only  → score × semantic_weight
#   • 'keyword' only   → score × keyword_weight
#   • 'both'           → score × (semantic_weight + keyword_weight)   ← bonus
#
# Reranking is skipped automatically when the candidate pool is too small,
# avoiding unnecessary compute overhead on short documents.

def rerank_contexts_batch(
    queries: List[str], contexts_per_query: List[List[Tuple[str, str]]], services: Dict[str, Any],
    semantic_weight: float, keyword_weight: float, k_rerank: int
) -> List[List[str]]:
    """
    Performs batched reranking with weighted scores across multiple queries.

    All (query, context) pairs from every query are scored in a single
    CrossEncoder batch call (batch_size=128) for GPU efficiency.  Scores are
    then adjusted by retrieval-source weights before the top-k chunks per
    query are selected and returned.
    """
    if not any(contexts_per_query): return [[] for _ in queries]
    t0 = time.perf_counter()
    all_pairs, all_sources = [], []
    contexts_counts = [len(contexts) for contexts in contexts_per_query]
    
    for query, contexts_with_sources in zip(queries, contexts_per_query):
        if contexts_with_sources:
            for context, source in contexts_with_sources:
                all_pairs.append([query, context])
                all_sources.append(source)
    if not all_pairs: return [[] for _ in queries]

    print(f" Reranking {len(all_pairs)} query-context pairs...")
    base_scores = services["reranker_model"].predict(all_pairs, show_progress_bar=False, batch_size=128)
    
    adjusted_scores = []
    for score, source in zip(base_scores, all_sources):
        if source == 'semantic': adjusted_scores.append(score * semantic_weight)
        elif source == 'keyword': adjusted_scores.append(score * keyword_weight)
        elif source == 'both': adjusted_scores.append(score * (semantic_weight + keyword_weight))
        else: adjusted_scores.append(score)

    final_reranked_contexts = []
    current_index = 0
    for count in contexts_counts:
        if count == 0:
            final_reranked_contexts.append([])
            continue
        query_pairs = all_pairs[current_index : current_index + count]
        query_adj_scores = adjusted_scores[current_index : current_index + count]
        original_contexts = [pair[1] for pair in query_pairs]
        reranked_results = sorted(zip(original_contexts, query_adj_scores), key=lambda x: x[1], reverse=True)
        final_chunks = [result[0] for result in reranked_results[:k_rerank]]
        final_reranked_contexts.append(final_chunks)
        current_index += count
            
    logtime(f"BATCH Weighted Reranking (top {k_rerank})", t0)
    return final_reranked_contexts

# --- REPLACE this entire function in worker.py ---

async def get_final_contexts_for_queries(
    queries: List[Tuple[str, str]], services: Dict[str, Any], semantic_weight: float,
    keyword_weight: float, k_semantic: int, k_keyword: int, k_rerank: int
) -> List[List[str]]:
    """
    Orchestrates retrieval and dynamically decides whether to rerank.

    Fires all hybrid_retrieve_async calls concurrently, then inspects the
    candidate pool for each query individually.  If a query has fewer
    candidates than k_rerank, the reranker is skipped (not enough signal
    to be useful).  Otherwise, normalised weights are passed to
    rerank_contexts_batch for precise final scoring.
    """
    retrieval_tasks = [hybrid_retrieve_async(query, doc_id, services, k_semantic, k_keyword) for query, doc_id in queries]
    all_candidate_contexts = await asyncio.gather(*retrieval_tasks)

    final_results = []
    # We check each query's results individually to decide if reranking is needed
    for i, candidate_contexts in enumerate(all_candidate_contexts):
        query, _ = queries[i]
        
        # If there are fewer candidates than k_rerank, reranking is just extra overhead.
        if len(candidate_contexts) < k_rerank:
            print(f"️ Skipping reranker for query '{query[:30]}...': not enough candidates ({len(candidate_contexts)}).")
            final_results.append(list(dict.fromkeys([ctx[0] for ctx in candidate_contexts])))
            continue

        # Otherwise, for larger documents, proceed with reranking
        print(f" Reranking {len(candidate_contexts)} candidates for query '{query[:30]}...'")
        total_weight = semantic_weight + keyword_weight
        norm_semantic_weight, norm_keyword_weight = (0.5, 0.5) if total_weight == 0 else (semantic_weight / total_weight, keyword_weight / total_weight)
        
        reranked_batch = rerank_contexts_batch(
            [query], [candidate_contexts], services,
            norm_semantic_weight, norm_keyword_weight, k_rerank
        )
        final_results.extend(reranked_batch)
        
    return final_results


# =============================================================================
# SECTION 13: LLM ANSWER GENERATION — THE FINAL STEP
# =============================================================================
# After retrieval and reranking, the top chunks are assembled into a context
# window and sent to the LLM with a carefully designed prompt.
#
# The prompt uses a Chain-of-Thought (CoT) structure:
#   Step 1 → LLM reasons internally over the context.
#   Step 2 → LLM writes a concise final answer from that reasoning.
#
# This is one of the most straightforward sections of the codebase:
# it's essentially "build a prompt, call the API, return the text."
# The helper get_llm_answers_in_batch fires all queries concurrently,
# and answer_queries ties the entire end-to-end pipeline together.

async def get_llm_answer_async(query: str, context_chunks: List[str], services: Dict[str, Any]) -> str:
    """
    Generates a final answer using the LLM with a Chain-of-Thought prompt.

    Assembles the top retrieved chunks into a single context string, injects
    them into a structured prompt alongside the user question, and calls
    the async OpenAI chat-completions API.  Returns a plain-text answer.
    """
    if not context_chunks:
        print(f"️ No context found for query: '{query}'. Returning default message.")
        return "Could not find relevant information in the document to answer this question."

    print(f"\n======================\n️  LLM Context for Question: {query}")
    log_context = "\n---\n".join(f"[Chunk {i+1}] {chunk[:250]}..." for i, chunk in enumerate(context_chunks))
    print(f"{log_context}\n======================\n")
    
    context = "\n\n---\n\n".join(context_chunks)
    prompt = f"""
    You are a highly precise, domain-agnostic professional assistant. Your task is to answer the user's question based *only* on the provided context.
    **Step 1: Internal Reasoning (Chain of Thought)**
    First, reason step-by-step to construct the answer based on the context.
    **Step 2: Final Answer Generation**
    Based on your reasoning, generate the final, concise answer as a single, clean paragraph.
    **CRITICAL RULES:**
    1.  **Strict Source Adherence**: Your answer must be explicitly supported by the context. If not, state that clearly.
    2.  **Language**: Answer in the language of the QUESTION.
    3.  **Special Case**: A query for a "secret token" requires you to provide only the token value.
    4.  When encountering puzzle type pdfs,estabish strong relationships between entitites and try to give a proper reasoning before giving your answer.The reasoning is for your understanding only.---
    **CONTEXT:**
    {context}
    ---
    **QUESTION:**
    {query}
    ---
    **FINAL ANSWER (as a single paragraph):**
    """
    try:
        start = time.perf_counter()
        response = await services["async_openai_client"].chat.completions.create(
            model="gpt-4o-mini", messages=[{"role": "user", "content": prompt}], temperature=0.0
        )
        logtime("LLM answer generation", start)
        return response.choices[0].message.content.strip()
    except Exception as e:
        print(f" Error getting LLM answer for query '{query}': {e}")
        return "An error occurred while generating the answer."

async def get_llm_answers_in_batch(tasks: List[Tuple[str, List[str]]], services: Dict[str, Any]) -> List[str]:
    """
    Gets LLM answers for all (query, context) tuples in parallel.

    Wraps each individual get_llm_answer_async call into a coroutine and
    fires them all concurrently with asyncio.gather, minimising total
    wall-clock time when processing multiple questions at once.
    """
    async_tasks = [get_llm_answer_async(query, context, services) for query, context in tasks]
    return await asyncio.gather(*async_tasks)

async def answer_queries(
    queries: List[Tuple[str, str]], services: Dict[str, Any], semantic_weight: float,
    keyword_weight: float, k_semantic: int, k_keyword: int, k_rerank: int
) -> List[str]:
    """
    Orchestrates the entire end-to-end RAG pipeline for a batch of queries.

    This is the single public function a caller needs: pass in a list of
    (query, document_id) tuples and receive back a list of plain-text answers
    in the same order.  Internally it chains retrieval → reranking → LLM.
    """
    final_contexts_per_query = await get_final_contexts_for_queries(
        queries=queries, services=services, semantic_weight=semantic_weight,
        keyword_weight=keyword_weight, k_semantic=k_semantic,
        k_keyword=k_keyword, k_rerank=k_rerank
    )
    llm_tasks = list(zip([q[0] for q in queries], final_contexts_per_query))
    return await get_llm_answers_in_batch(llm_tasks, services)
