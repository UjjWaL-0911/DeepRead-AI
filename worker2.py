# --- SECTION 1: IMPORTS & INITIAL SETUP ---

import os
import asyncio
import hashlib
import json
import re
import time
from collections import defaultdict
from typing import List, Tuple, Dict, Any
from urllib.parse import urlparse, parse_qsl, urlencode
from pptx import Presentation
import base64
import pandas as pd

# Third-party libraries
import aiohttp
import fitz  # PyMuPDF
import nltk
import openai
import redis
from bs4 import BeautifulSoup
from docx import Document
from dotenv import load_dotenv
from langchain.text_splitter import RecursiveCharacterTextSplitter
from sentence_transformers import CrossEncoder
from astrapy import DataAPIClient
from langdetect import detect
from deep_translator import GoogleTranslator

# Perform initial setup tasks
print("🚀 Worker starting up...")
load_dotenv()
nltk.download('punkt', quiet=True)
nltk.download('stopwords', quiet=True)
os.environ["TOKENIZERS_PARALLELISM"] = "false"


# --- SECTION 2: CONFIGURATION & INITIALIZATION ---

# --- Constants ---
EMBEDDING_MODEL_API = "text-embedding-3-small"
EMBEDDING_DIMENSION = 1536
VECTOR_COLLECTION_NAME = "newtest"
KEYWORD_INDEX_COLLECTION_NAME = "keyword_index"
PROCESSED_DOCS_KEY = "v1" # Redis key for tracking processed docs

def load_environment_config() -> Dict[str, str]:
    """Loads required environment variables and returns them as a dictionary."""
    config = {
        "redis_url": os.getenv("REDIS_URL"),
        "openai_api_key": os.getenv("OPENAI_API_KEY"),
        "astra_db_api_endpoint": os.getenv("ASTRA_DB_API_ENDPOINT"),
        "astra_db_application_token": os.getenv("ASTRA_DB_APPLICATION_TOKEN"),
    }
    if not all(config.values()):
        raise ValueError("One or more essential environment variables are missing.")
    return config

def initialize_services(config: Dict[str, str]) -> Dict[str, Any]:
    """Initializes and returns clients for external services and ML models."""
    print("⚙️ Initializing external services and models...")

    # Initialize clients
    redis_conn = redis.from_url(config["redis_url"])
    sync_openai_client = openai.OpenAI(api_key=config["openai_api_key"])
    async_openai_client = openai.AsyncOpenAI(api_key=config["openai_api_key"])
    
    print("🧠 Loading CrossEncoder model for reranking...")
    reranker_model = CrossEncoder('BAAI/bge-reranker-base')
    print("✅ CrossEncoder model loaded.")
    
    print("☁️ Connecting to Astra DB...")
    astra_client = DataAPIClient(config["astra_db_application_token"])
    db = astra_client.get_database(config["astra_db_api_endpoint"])
    print(f"✅ Connected to Astra DB. Namespace: {db.namespace}")

    # Setup Astra DB collections
    collection_names = db.list_collection_names()
    
    # Vector collection
    if VECTOR_COLLECTION_NAME in collection_names:
        print(f"Collection '{VECTOR_COLLECTION_NAME}' already exists. Loading it.")
        vector_collection = db.get_collection(VECTOR_COLLECTION_NAME)
    else:
        print(f"Collection '{VECTOR_COLLECTION_NAME}' not found. Creating it...")
        vector_collection = db.create_collection(
            VECTOR_COLLECTION_NAME,
            dimension=EMBEDDING_DIMENSION
        )
        print("✅ Vector collection created.")

    # Keyword index collection
    if KEYWORD_INDEX_COLLECTION_NAME in collection_names:
        print(f"Collection '{KEYWORD_INDEX_COLLECTION_NAME}' already exists. Loading it.")
        keyword_index_collection = db.get_collection(KEYWORD_INDEX_COLLECTION_NAME)
    else:
        print(f"Collection '{KEYWORD_INDEX_COLLECTION_NAME}' not found. Creating it...")
        keyword_index_collection = db.create_collection(KEYWORD_INDEX_COLLECTION_NAME)
        print("✅ Keyword index collection created.")

    print("✅ Astra DB Collections ready.")
    print("-" * 50)
    
    return {
        "redis_conn": redis_conn,
        "sync_openai_client": sync_openai_client,
        "async_openai_client": async_openai_client,
        "reranker_model": reranker_model,
        "db": db,
        "vector_collection": vector_collection,
        "keyword_index_collection": keyword_index_collection
    }

# --- RAG Pipeline Configuration ---
# Central place for all tunable parameters of the RAG pipeline.
RAG_CONFIG = {
    "semantic_weight": 0.7,
    "keyword_weight": 0.3,
    "k_semantic": 30,
    "k_keyword": 30,
    "k_rerank": 15
}
print(f"✅ RAG pipeline configured with default tuning parameters: {RAG_CONFIG}")


# --- SECTION 3: UTILITY FUNCTIONS ---

def logtime(phase: str, start: float, last: float = None) -> float:
    """Logs the time elapsed since the start and the last checkpoint."""
    now = time.perf_counter()
    since_start = now - start
    log_message = f"⏱️ [{phase}] {since_start:.2f}s"
    if last:
        since_last = now - last
        log_message += f" (+{since_last:.2f}s since last)"
    print(log_message)
    return now

def generate_document_id(source: str) -> str:
    """
    Generates a stable, unique document ID from a URL or file path.
    Hashes the final ID to prevent excessive length.
    """
    if urlparse(source).scheme in ['http', 'https']:
        parsed = urlparse(source)
        base_id = parsed.path.strip("/").replace("/", "-") or "root"
        if parsed.query:
            query_string = urlencode(sorted(parse_qsl(parsed.query)))
            full_id = f"{base_id}?{query_string}"
        else:
            full_id = base_id
    else: # It's a file path
        full_id = os.path.basename(source)
    
    hashed = hashlib.md5(full_id.encode()).hexdigest()[:8]
    sanitized_base = re.sub(r'[^a-zA-Z0-9_-]', '', os.path.splitext(full_id)[0])
    return f"{sanitized_base}-{hashed}"


# --- SECTION 4: DATA EXTRACTION & PREPARATION ---

def extract_text_from_pdf(file_path: str) -> Tuple[str, int]:
    """Extracts all text and page count from a PDF file."""
    doc = fitz.open(file_path)
    full_text = "".join(page.get_text() for page in doc)
    num_pages = doc.page_count
    doc.close()
    return full_text, num_pages

def extract_text_from_pptx(file_path: str) -> Tuple[str, int]:
    """Extracts all text and slide count from a PPTX file."""
    prs = Presentation(file_path)
    full_text = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                full_text.append(shape.text)
    return "\n".join(full_text), len(prs.slides)

def extract_text_from_docx(file_path: str) -> Tuple[str, int]:
    """Extracts all text and estimates page count from a DOCX file."""
    doc = Document(file_path)
    full_text = "\n".join([para.text for para in doc.paragraphs if para.text.strip()])
    para_count = len(doc.paragraphs)
    num_pages = max(1, para_count // 5)  # Heuristic for page count
    return full_text, num_pages

def extract_text_from_html(html_content: str) -> Tuple[str, int]:
    """Extracts clean text and estimates page count from HTML content."""
    soup = BeautifulSoup(html_content, "html.parser")
    for tag in soup(["script", "style", "noscript", "header", "footer", "nav"]):
        tag.extract()
    text = soup.get_text(separator="\n", strip=True)
    num_pages = max(1, len(text.split()) // 300) # Heuristic for page count
    return text, num_pages

async def scrape_and_clean_url(session: aiohttp.ClientSession, url: str) -> str:
    """Asynchronously scrapes and cleans text content from a given URL."""
    try:
        async with session.get(url, timeout=15) as response:
            response.raise_for_status()
            html = await response.text()
            clean_text, _ = extract_text_from_html(html)
            return clean_text
    except Exception as e:
        print(f"⚠️ Could not scrape URL {url}: {e}")
        return ""
    
def extract_text_from_excel(file_path: str) -> Tuple[str, int]:
    """Extracts text from all sheets of an Excel file."""
    xls = pd.ExcelFile(file_path)
    full_text = []
    for sheet_name in xls.sheet_names:
        df = pd.read_excel(xls, sheet_name=sheet_name, header=None)
        # Add a header for context and convert sheet to string
        sheet_text = f"--- Sheet: {sheet_name} ---\n{df.to_string(index=False)}"
        full_text.append(sheet_text)
    return "\n\n".join(full_text), len(xls.sheet_names)

async def extract_text_from_image_async(file_path: str, services: Dict[str, Any]) -> Tuple[str, int]:
    """Extracts text from an image using OpenAI's vision model."""
    print("👁️ Extracting text from image using Vision API...")
    with open(file_path, "rb") as image_file:
        base64_image = base64.b64encode(image_file.read()).decode('utf-8')

    response = await services["async_openai_client"].chat.completions.create(
        model="gpt-4o-mini",
        messages=[
            {
                "role": "user",
                "content": [
                    {"type": "text", "text": "Extract all text from this image exactly as it appears. If no text is present, say nothing."},
                    {"type": "image_url", "image_url": {"url": f"data:image/jpeg;base64,{base64_image}"}},
                ],
            }
        ],
        max_tokens=2048,
    )
    extracted_text = response.choices[0].message.content
    return extracted_text, 1 # An image is considered a single page


# --- SECTION 5: CORE INDEXING COMPONENTS ---
import nltk
from nltk.corpus import stopwords

# Create a set of standard English stopwords from the NLTK library.
# The try/except block will automatically download the list if it's not present.
try:
    NLTK_STOPWORDS = set(stopwords.words('english'))
except LookupError:
    print("Downloading NLTK stopwords list...")
    nltk.download('stopwords')
    NLTK_STOPWORDS = set(stopwords.words('english'))

print(f"✅ Standard NLTK stopword list initialized with {len(NLTK_STOPWORDS)} words.")



def simple_doc_type_classifier(text: str) -> str:
    """Classifies document type based on keywords for chunking strategy."""
    lowered = text.lower()
    if any(word in lowered for word in ["policy", "contract", "agreement", "legal"]):
        return "policy"
    elif any(word in lowered for word in ["research", "study", "technical", "scientific"]):
        return "scientific"
    elif any(word in lowered for word in ["story", "novel", "narrative", "literature"]):
        return "narrative"
    return "general"

def get_chunking_params(doc_type: str, num_pages: int) -> Tuple[int, int]:
    """Determines chunk size and overlap based on document type and length."""
    # This logic is kept as per the original script. It can be expanded.
    chunk_size = 1000
    chunk_overlap = 100
    # print(f"ℹ️ Using chunking params: size={chunk_size}, overlap={chunk_overlap} for type='{doc_type}'")
    return chunk_size, chunk_overlap

def chunk_text(text: str, doc_type: str, num_pages: int) -> List[str]:
    """Splits text into chunks using a recursive character splitter."""
    chunk_size, chunk_overlap = get_chunking_params(doc_type, num_pages)
    splitter = RecursiveCharacterTextSplitter(
        chunk_size=chunk_size, 
        chunk_overlap=chunk_overlap
    )
    return splitter.split_text(text)

def get_embeddings_sync(texts: List[str], client: openai.OpenAI) -> List[List[float]]:
    """Synchronously generates embeddings for a list of texts."""
    response = client.embeddings.create(
        input=[t.replace('\n', ' ') for t in texts], 
        model=EMBEDDING_MODEL_API
    )
    return [e.embedding for e in response.data]

async def get_embeddings_batch_async(
    texts: List[str], client: openai.AsyncOpenAI, batch_size: int = 500
) -> List[List[float]]:
    """Generates embeddings for a list of texts in batches asynchronously."""
    all_embeddings = []
    tasks = []
    for i in range(0, len(texts), batch_size):
        batch = [t.replace('\n', ' ') for t in texts[i:i + batch_size]]
        tasks.append(client.embeddings.create(input=batch, model=EMBEDDING_MODEL_API))
    
    responses = await asyncio.gather(*tasks)
    for response in responses:
        all_embeddings.extend([e.embedding for e in response.data])
    return all_embeddings


def tokenize_for_indexing(text: str) -> List[str]:
    """
    Lowercases, tokenizes text, and removes standard NLTK stopwords for
    efficient and relevant keyword indexing.
    """
    tokens = re.findall(r'\w+', text.lower())
    # Filter out any token that is in our NLTK_STOPWORDS set
    return [token for token in tokens if token not in NLTK_STOPWORDS]


def build_and_store_inverted_index(
    document_id: str, chunks: List[str], keyword_collection
):
    """Builds an inverted index and stores it in Astra DB."""
    inverted_index = defaultdict(set)
    for i, chunk in enumerate(chunks):
        chunk_id = f"{document_id}-chunk-{i}"
        for token in tokenize_for_indexing(chunk):
            inverted_index[token].add(chunk_id)

    index_docs = [
        {
            "_id": f"{document_id}-{token}",
            "document_id": document_id,
            "token": token,
            "chunks": list(chunk_ids)
        }
        for token, chunk_ids in inverted_index.items()
    ]
    
    if index_docs:
        keyword_collection.insert_many(index_docs)
    print(f"✅ Inverted index stored for {len(inverted_index)} unique tokens.")


# --- SECTION 6: FILE PROCESSING PIPELINES (with Redis check) ---

# In worker2.0.py, MODIFY the '_process_and_index_document' function

async def _process_and_index_document(
    document_id: str,
    text: str,
    num_pages: int,
    services: Dict[str, Any]
):
    """
    Internal function now with parallelized indexing and embedding.
    """
    t0 = time.perf_counter()
    
    doc_type = simple_doc_type_classifier(text)
    chunks = chunk_text(text, doc_type, num_pages)
    if not chunks:
        print(f"⚠️ No chunks could be generated for '{document_id}'. Aborting.")
        return
    t1 = logtime(f"Split text into {len(chunks)} chunks", t0)

    # --- Run Indexing and Embedding in Parallel ---
    print("🚀 Starting parallel embedding generation and keyword indexing...")
    embedding_task = get_embeddings_batch_async(chunks, services["async_openai_client"])
    keyword_task = asyncio.to_thread(
        build_and_store_inverted_index,
        document_id, chunks, services["keyword_index_collection"]
    )
    
    dense_embeddings, _ = await asyncio.gather(embedding_task, keyword_task)
    t2 = logtime("Finished parallel embedding and indexing", t0, t1)
    
    # Insert documents into vector collection
    documents_to_insert = [
        {
            "_id": f"{document_id}-chunk-{i}", "document_id": document_id,
            "text": chunk, "$vector": dense,
        }
        for i, (chunk, dense) in enumerate(zip(chunks, dense_embeddings))
    ]
    services["vector_collection"].insert_many(documents_to_insert)
    t3 = logtime("Inserted all chunks into Astra DB", t0, t2)
    
    services["redis_conn"].sadd(PROCESSED_DOCS_KEY, document_id)
    print(f"📝 Marked '{document_id}' as processed in Redis.")
    
    print(f"✅ Finished indexing '{document_id}' successfully.")
    print("-" * 50)
    
    
async def translate_if_needed_async(text: str, target_lang: str = "en") -> str:
    """Detects language and translates it to the target language if different, using the deep-translator library."""
    try:
        # Step 1: Detect the language of the input text.
        detected_lang = await asyncio.to_thread(detect, text)

        # Step 2: Check if translation is necessary.
        if detected_lang != target_lang:
            print(f"🌐 Detected language: {detected_lang}, translating to {target_lang}...")

            # Step 3: Perform the translation using the correct deep-translator syntax.
            translated_text = await asyncio.to_thread(
                GoogleTranslator(source='auto', target=target_lang).translate,
                text
            )

            # Step 4: Return the formatted, translated text.
            if translated_text:
                return f"\n\n--- AUTO TRANSLATION ({detected_lang} → {target_lang}) ---\n{translated_text}\n--- END TRANSLATION ---"
            else:
                return ""  # Return empty string if translation result is empty.
        else:
            # The text is already in the target language.
            return ""

    except Exception as e:
        print(f"⚠️ Language detection/translation failed: {e}")
        return ""


async def process_and_index_pdf(file_path: str, document_id: str, services: Dict[str, Any]):
    """Orchestrates the full processing pipeline for a PDF file, using a PRE-GENERATED document_id."""
    if services["redis_conn"].sismember(PROCESSED_DOCS_KEY, document_id):
        print(f"ℹ️ Document ID '{document_id}' (from PDF) already processed. Skipping.")
        return

    print(f"⚙️ Starting pipeline for PDF: {file_path} (ID: {document_id})")
    full_text, num_pages = extract_text_from_pdf(file_path)
    if not full_text.strip():
        print(f"⚠️ No text could be extracted from the PDF '{file_path}'. Aborting.")
        return

    found_urls = sorted(list(set(re.findall(r'https?://[^\s<>"]+|www\.[^\s<>"]+', full_text))))
    combined_text = full_text

    if found_urls:
        print(f"🔗 Found {len(found_urls)} unique URLs. Scraping...")
        async with aiohttp.ClientSession() as session:
            scraping_tasks = [scrape_and_clean_url(session, url) for url in found_urls]
            scraped_results = await asyncio.gather(*scraping_tasks)
            scraped_texts = [
                f"--- LINKED CONTENT FROM: {url} ---\n{text}\n--- END LINKED CONTENT ---)"
                for url, text in zip(found_urls, scraped_results) if text
            ]
            if scraped_texts:
                combined_text += "\n\n" + "\n\n".join(scraped_texts)
                print(f"✅ Appended content from {len(scraped_texts)} URLs.")

    # 🔹 Automatic language detection & translation
    translation_context = await translate_if_needed_async(combined_text)
    if translation_context:
        combined_text += translation_context

    await _process_and_index_document(document_id, combined_text, num_pages, services)


async def process_and_index_docx(file_path: str, document_id: str, services: Dict[str, Any]):
    """Orchestrates the full processing pipeline for a DOCX file, using a PRE-GENERATED document_id."""
    if services["redis_conn"].sismember(PROCESSED_DOCS_KEY, document_id):
        print(f"ℹ️ Document ID '{document_id}' (from DOCX) already processed. Skipping.")
        return

    print(f"⚙️ Starting pipeline for DOCX: {file_path} (ID: {document_id})")
    full_text, num_pages = extract_text_from_docx(file_path)
    if not full_text.strip():
        print(f"⚠️ No text could be extracted from the DOCX '{file_path}'. Aborting.")
        return
    await _process_and_index_document(document_id, full_text, num_pages, services)


async def process_and_index_url(url: str, document_id: str, services: Dict[str, Any]):
    """Orchestrates the full processing pipeline for a web URL, using a PRE-GENERATED document_id."""
    if services["redis_conn"].sismember(PROCESSED_DOCS_KEY, document_id):
        print(f"ℹ️ Document ID '{document_id}' (from URL) already processed. Skipping.")
        return

    print(f"⚙️ Starting pipeline for URL: {url} (ID: {document_id})")
    async with aiohttp.ClientSession() as session:
        html_content = await scrape_and_clean_url(session, url)
    if not html_content:
        print(f"⚠️ Could not retrieve or parse content from URL '{url}'. Aborting.")
        return
        
    full_text, num_pages = extract_text_from_html(html_content)
    if not full_text.strip():
        print(f"⚠️ No text content found after cleaning HTML from '{url}'. Aborting.")
        return
    await _process_and_index_document(document_id, full_text, num_pages, services)


async def process_and_index_pptx(file_path: str, document_id: str, services: Dict[str, Any]):
    """Orchestrates the full processing pipeline for a PPTX file."""
    if services["redis_conn"].sismember(PROCESSED_DOCS_KEY, document_id):
        print(f"ℹ️ Document ID '{document_id}' (from PPTX) already processed. Skipping.")
        return
    
    print(f"⚙️ Starting pipeline for PPTX: {file_path} (ID: {document_id})")
    full_text, num_pages = extract_text_from_pptx(file_path)
    if not full_text.strip():
        print(f"⚠️ No text could be extracted from the PPTX '{file_path}'. Aborting.")
        return
    await _process_and_index_document(document_id, full_text, num_pages, services)


async def process_and_index_excel(file_path: str, document_id: str, services: Dict[str, Any]):
    """Orchestrates the full processing pipeline for an Excel file."""
    if services["redis_conn"].sismember(PROCESSED_DOCS_KEY, document_id):
        print(f"ℹ️ Document ID '{document_id}' (from Excel) already processed. Skipping.")
        return
        
    print(f"⚙️ Starting pipeline for Excel: {file_path} (ID: {document_id})")
    full_text, num_sheets = extract_text_from_excel(file_path)
    if not full_text.strip():
        print(f"⚠️ No text could be extracted from the Excel file '{file_path}'. Aborting.")
        return
    await _process_and_index_document(document_id, full_text, num_sheets, services)

async def process_and_index_image(file_path: str, document_id: str, services: Dict[str, Any]):
    """Orchestrates the full processing pipeline for an image file."""
    if services["redis_conn"].sismember(PROCESSED_DOCS_KEY, document_id):
        print(f"ℹ️ Document ID '{document_id}' (from Image) already processed. Skipping.")
        return

    print(f"⚙️ Starting pipeline for Image: {file_path} (ID: {document_id})")
    full_text, num_pages = await extract_text_from_image_async(file_path, services)
    if not full_text.strip():
        print(f"⚠️ No text could be extracted from the Image file '{file_path}'. Aborting.")
        return
    await _process_and_index_document(document_id, full_text, num_pages, services)

async def process_and_index_zip(file_path: str, document_id: str, services: Dict[str, Any]):
    """Orchestrates the processing of a ZIP archive by processing its contents."""
    if services["redis_conn"].sismember(PROCESSED_DOCS_KEY, document_id):
        print(f"ℹ️ Document ID '{document_id}' (from ZIP) already processed. Skipping.")
        return

    print(f"⚙️ Starting pipeline for ZIP: {file_path} (ID: {document_id})")
    all_texts = []
    total_pages = 0

    with tempfile.TemporaryDirectory() as temp_dir:
        with zipfile.ZipFile(file_path, 'r') as zip_ref:
            zip_ref.extractall(temp_dir)
        
        print(f"🗃️ Extracted {len(os.listdir(temp_dir))} files from zip archive.")
        for filename in os.listdir(temp_dir):
            extracted_file_path = os.path.join(temp_dir, filename)
            file_ext = os.path.splitext(filename)[-1].lower()
            text, pages = "", 0
            
            try:
                print(f"   -> Processing '{filename}'...")
                if file_ext == '.pdf':
                    text, pages = extract_text_from_pdf(extracted_file_path)
                elif file_ext == '.docx':
                    text, pages = extract_text_from_docx(extracted_file_path)
                elif file_ext == '.pptx':
                    text, pages = extract_text_from_pptx(extracted_file_path)
                elif file_ext in ['.xlsx', '.xls']:
                    text, pages = extract_text_from_excel(extracted_file_path)
                elif file_ext in ['.jpg', '.jpeg', '.png']:
                    # Image extraction is async, so we await it here.
                    text, pages = await extract_text_from_image_async(extracted_file_path, services)
                else:
                    print(f"   -> ⚠️ Unsupported file type '{file_ext}' in ZIP. Skipping.")
                    continue
                
                if text.strip():
                    all_texts.append(f"--- START CONTENT FROM: {filename} ---\n\n{text}\n\n--- END CONTENT FROM: {filename} ---")
                    total_pages += pages
            except Exception as e:
                print(f"   -> ❌ Error processing '{filename}' inside ZIP: {e}")

    if not all_texts:
        print(f"⚠️ No text could be extracted from any files in the ZIP '{file_path}'. Aborting.")
        return

    combined_text = "\n\n".join(all_texts)
    await _process_and_index_document(document_id, combined_text, total_pages, services)








## --- SECTION 7: HYBRID RETRIEVAL & RERANKING LOGIC ---
# ADD THIS NEW FUNCTION (can go in Section 7 or a new Section 9)
def get_all_chunks_for_document(document_id: str, services: Dict[str, Any]) -> List[Dict]:
    """Retrieves all chunk documents from Astra DB for a given document_id."""
    print(f"Retrieving all chunks for document: {document_id}")
    # Use a filter to find all chunks belonging to the document
    results = services["vector_collection"].find(filter={"document_id": document_id})
    return list(results)

# ADD THIS NEW FUNCTION to worker2.0.py

async def discover_concepts_dynamically_async(document_id: str, services: Dict[str, Any]) -> List[str]:
    """
    Phase 0: Scans a sample of the document to have the LLM dynamically identify
    the key concepts that should be tracked for consistency.
    """
    print("🤖 Starting Phase 0: Dynamic Concept Discovery...")
    
    # Define a default list to fall back on in case of failure
    fallback_concepts = [
        "Principal Amount", "Interest Rate", "Effective Date", "Maturity Date",
        "Closing Date", "Governing Law", "Lender Name", "Borrower Name"
    ]

    # Get a sample of chunks (e.g., first 3 and last 2) to get a feel for the document
    all_chunks = await asyncio.to_thread(get_all_chunks_for_document, document_id, services)
    if not all_chunks:
        print("⚠️ No chunks found for concept discovery. Using fallback list.")
        return fallback_concepts

    sample_chunks = all_chunks[:3] + all_chunks[-2:]
    sample_text = "\n---\n".join([chunk['text'] for chunk in sample_chunks])

    # This is a new, specialized prompt for concept discovery
    prompt = f"""
    You are a senior analyst and document specialist. Your goal is to identify the most critical and recurring concepts in a document that should be tracked for consistency and summarization.

    **Instructions:**
    - Analyze the provided sample text from a larger document.
    - Identify between 5 and 15 key trackable concepts.
    - Focus on specific entities, dates, monetary values, and key terms that define obligations or agreements (e.g., 'Loan Maturity Date' is better than just 'Date').
    - The concepts should be the "column headers" you would use if you were summarizing this document in a spreadsheet.

    **Example:**
    - If the text is a loan agreement, concepts might be ["Principal Loan Amount", "Interest Rate", "Lender Name", "Borrower Name"].
    - If the text is a scientific paper, concepts might be ["Hypothesis", "Sample Size", "Key Finding", "Primary Author", "Publication Date"].

    **Your Task:**
    Based on the sample text below, generate the list of key concepts.
    Return your answer ONLY as a valid JSON list of strings.

    **Sample Text:**
    "{sample_text}"
    """

    try:
        discovered_concepts = await get_llm_json_response_async(prompt, services)
        if isinstance(discovered_concepts, list) and len(discovered_concepts) > 0:
            print(f"✅ Dynamic concepts discovered: {discovered_concepts}")
            return discovered_concepts
        else:
            print("⚠️ Dynamic discovery did not return a valid list. Using fallback.")
            return fallback_concepts
    except Exception as e:
        print(f"❌ Error during concept discovery: {e}. Using fallback list.")
        return fallback_concepts
    
    
    

# ADD THIS NEW FUNCTION
from collections import defaultdict

async def extract_all_facts_from_document(document_id: str, services: Dict[str, Any], concepts_to_track: List[str]) -> Dict:
    """
    Phase 1: Iterates through document chunks IN BATCHES to extract a structured 'Fact Sheet' much faster.
    """
    print(f"🤖 Starting Phase 1: Batched Fact Extraction using concepts: {concepts_to_track}")
    
    all_chunks = await asyncio.to_thread(get_all_chunks_for_document, document_id, services)
    aggregated_facts = defaultdict(list)
    
    # --- BATCHING LOGIC ---
    BATCH_SIZE = 5  # Process 5 chunks per API call. Adjust based on chunk size and context limits.
    chunk_batches = [all_chunks[i:i + BATCH_SIZE] for i in range(0, len(all_chunks), BATCH_SIZE)]
    print(f"Split {len(all_chunks)} chunks into {len(chunk_batches)} batches of up to {BATCH_SIZE} chunks each.")

    tasks = []
    for batch in chunk_batches:
        # --- NEW MEGA-PROMPT for BATCHING ---
        # Create a formatted string containing all chunks in the batch, each with a unique ID
        formatted_batch_text = "\n\n---\n\n".join(
            f"--- CHUNK ID: {chunk['_id']} ---\n{chunk['text']}" for chunk in batch
        )

        prompt = f"""
        You are an expert legal data extraction bot. Your task is to meticulously analyze a batch of text chunks from a single document.

        **Instructions:**
        - I will provide several text chunks, each marked with a unique "CHUNK ID".
        - For EACH chunk, extract values for the concepts: {concepts_to_track}.
        - Apply these normalization rules:
            - **dates**: Normalize to `YYYY-MM-DD` format.
            - **monetary values**: Extract only the numerical value (e.g., '$1,500,000' becomes `1500000`).
        - The value for each concept must be a list to handle multiple mentions within a single chunk.

        **Output Format:**
        Your entire response must be a single, valid JSON object. The top-level keys must be the exact "CHUNK ID"s from the input. The value for each CHUNK ID should be another JSON object containing the extracted concepts and their values for that chunk. If a chunk contains no relevant concepts, its value should be an empty object {{}}.

        **Example Output:**
        {{
          "doc1-chunk-5": {{
            "Effective Date": ["2025-08-29"],
            "Principal Amount": [1000000]
          }},
          "doc1-chunk-6": {{}},
          "doc1-chunk-7": {{
            "Lender Name": ["Future Bank"]
          }}
        }}

        **Chunks to Analyze:**
        {formatted_batch_text}
        """
        tasks.append(get_llm_json_response_async(prompt, services))

    # Run all batch tasks concurrently
    batch_responses = await asyncio.gather(*tasks)

    # --- PROCESS BATCH RESPONSES ---
    for response_data in batch_responses:
        if not isinstance(response_data, dict): continue

        for chunk_id, extracted_data in response_data.items():
            # Find the original source text for this chunk_id
            original_chunk_text = next((chunk['text'] for chunk in all_chunks if chunk['_id'] == chunk_id), "")
            
            for concept, values in extracted_data.items():
                if concept in concepts_to_track and isinstance(values, list):
                    for value in values:
                        aggregated_facts[concept].append({
                            "value": value,
                            "source_text": original_chunk_text
                        })
    
    print(f"✅ Batched fact extraction complete. Found data for {len(aggregated_facts)} concepts.")
    return dict(aggregated_facts)


# ADD THIS NEW FUNCTION
def analyze_facts_for_discrepancies(facts: Dict) -> List[Dict]:
    """
    Phase 2: Analyzes the structured 'Fact Sheet' to find inconsistencies.
    """
    findings = []
    for concept, values_list in facts.items():
        if len(values_list) > 1:
            # Normalize values for a more robust comparison.
            # This is a simple example; you can make this more sophisticated.
            unique_values = {str(item['value']).strip().lower() for item in values_list}

            if len(unique_values) > 1:
                finding = {
                    "type": "Data Inconsistency",
                    "severity": "High",
                    "message": f"Found multiple conflicting values for the concept: '{concept}'.",
                    "evidence": [item for item in values_list]
                }
                findings.append(finding)

    print(f"✅ Analysis complete. Found {len(findings)} potential discrepancies.")
    return findings

# ADD THIS NEW, SPECIALIZED FUNCTION to worker2.0.py
async def get_llm_json_response_async(prompt: str, services: Dict[str, Any]) -> Dict:
    """
    Calls the LLM with a prompt expecting a JSON response and handles parsing.
    """
    try:
        response = await services["async_openai_client"].chat.completions.create(
            model="gpt-4o-mini",
            messages=[{"role": "user", "content": prompt}],
            temperature=0.0,
            response_format={"type": "json_object"}, # Use JSON mode for reliability
        )
        content = response.choices[0].message.content
        return json.loads(content)
    except json.JSONDecodeError:
        print(f"⚠️ Warning: LLM did not return valid JSON. Response: {content}")
        return {}
    except Exception as e:
        print(f"❌ Error getting LLM JSON response: {e}")
        return {}

